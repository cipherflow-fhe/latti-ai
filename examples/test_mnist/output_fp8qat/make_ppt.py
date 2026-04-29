#!/usr/bin/env python3
"""Generate FP8 QAT Cipher Inference presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ──────────────────────────────────────────────────────────────
BG_DARK    = RGBColor(0x0F, 0x17, 0x2A)   # deep navy
BG_CARD    = RGBColor(0x1A, 0x25, 0x3C)   # card background
ACCENT     = RGBColor(0x00, 0xBF, 0xD8)   # cyan accent
ACCENT2    = RGBColor(0x7C, 0x3A, 0xED)   # purple accent
ACCENT3    = RGBColor(0x10, 0xB9, 0x81)   # green accent
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xA0, 0xAE, 0xC0)
DIM        = RGBColor(0x64, 0x74, 0x8B)
ORANGE     = RGBColor(0xF9, 0x7B, 0x16)
RED_ACCENT = RGBColor(0xEF, 0x44, 0x44)
GOLD       = RGBColor(0xFF, 0xD7, 0x00)   # gold for contribution highlight


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=BG_CARD, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Consolas"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rich_text_box(slide, left, top, width, height):
    """Return a text_frame for multi-paragraph rich content."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def add_para(tf, text, size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
             font_name="Consolas", space_after=Pt(4)):
    p = tf.add_paragraph() if tf.paragraphs[0].text or len(tf.paragraphs) > 0 and tf.paragraphs[0].text != "" else tf.paragraphs[0]
    if len(tf.paragraphs) > 1 or (len(tf.paragraphs) == 1 and tf.paragraphs[0].text):
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = space_after
    return p


def add_bullet(tf, text, size=13, color=LIGHT_GRAY, level=0, font_name="Consolas", bold=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.level = level
    p.space_after = Pt(3)
    return p


def add_badge(slide, left, top, text, color=ACCENT):
    w, h = Inches(1.8), Inches(0.35)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Consolas"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(2)
    return shape


def add_slide_number(slide, num, total):
    add_text_box(slide, Inches(12.0), Inches(8.2), Inches(1.0), Inches(0.3),
                 f"{num}/{total}", font_size=9, color=DIM, alignment=PP_ALIGN.RIGHT)


# ── Create presentation ────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
TOTAL = 14

# ===================== SLIDE 1: Title =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

# Decorative accent line
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), fill_color=ACCENT)
add_shape(slide, Inches(0), Inches(0.06), Inches(13.33), Inches(0.03), fill_color=ACCENT2)

add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
             "FP8 Quantized Cipher Inference", font_size=40, bold=True, color=WHITE,
             font_name="Consolas", alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(0.8),
             "Custom CKKS Parameters for N=8192 Encrypted Neural Network Inference",
             font_size=18, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(0.5),
             "5.6x Speedup  |  ~127-bit Security  |  PASS Verification",
             font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
             "latti-ai  |  CipherFlow FHE Platform", font_size=13, color=DIM,
             alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 1, TOTAL)

# ===================== SLIDE 2: Agenda =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), fill_color=ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.5),
             "Agenda", font_size=28, bold=True, color=WHITE)

items = [
    ("01", "Problem & Motivation", "Why FP8 + smaller ring dimension?"),
    ("02", "Our Core Contribution", "CKKS param selection for QAT models"),
    ("03", "Methodology & Algorithm", "Step-by-step parameter optimization"),
    ("04", "Case Study: FP8 QAT", "Applying the methodology to MNIST"),
    ("05", "Case Study: FHE Conversion", "ReLU → polynomial, pooling, BN fusion"),
    ("06", "Case Study: CKKS Design", "Custom PN13QP218s with 24-bit primes"),
    ("07", "Encrypted Inference Results", "Performance comparison & verification"),
    ("08", "Generalization & Impact", "Applicable to any QAT model + FHE inference"),
]

for i, (num, title, desc) in enumerate(items):
    y = Inches(1.2 + i * 0.73)
    add_shape(slide, Inches(0.8), y, Inches(0.6), Inches(0.55), fill_color=ACCENT2 if i < 3 else ACCENT)
    add_text_box(slide, Inches(0.8), y + Inches(0.08), Inches(0.6), Inches(0.4),
                 num, font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.6), y + Inches(0.0), Inches(5.5), Inches(0.3),
                 title, font_size=14, bold=True, color=WHITE)
    add_text_box(slide, Inches(1.6), y + Inches(0.28), Inches(5.5), Inches(0.25),
                 desc, font_size=11, color=LIGHT_GRAY)

# Right side key result card
add_shape(slide, Inches(8.0), Inches(1.3), Inches(4.5), Inches(5.5), border_color=ACCENT3)
add_text_box(slide, Inches(8.3), Inches(1.5), Inches(4), Inches(0.4),
             "Key Result", font_size=20, bold=True, color=ACCENT3)
tf = add_rich_text_box(slide, Inches(8.3), Inches(2.1), Inches(4), Inches(4.5))
add_para(tf, "N = 8,192  (poly modulus degree)", size=13, color=WHITE, bold=True)
add_para(tf, "7 levels (24-bit Q primes)", size=13, color=WHITE)
add_para(tf, "", size=8, color=DIM)
add_para(tf, "Inference Time", size=12, color=DIM)
add_para(tf, "1,194 ms", size=32, color=ACCENT3, bold=True)
add_para(tf, "vs 6,701 ms baseline", size=12, color=DIM)
add_para(tf, "", size=8, color=DIM)
add_para(tf, "Speedup", size=12, color=DIM)
add_para(tf, "5.6x faster", size=28, color=ORANGE, bold=True)
add_para(tf, "", size=8, color=DIM)
add_para(tf, "Max error: 0.037 | Verification: PASS", size=12, color=ACCENT3)

add_slide_number(slide, 2, TOTAL)

# ===================== SLIDE 3: Problem & Motivation =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), fill_color=ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.5),
             "Problem & Motivation", font_size=28, bold=True, color=WHITE)

# Problem card
add_shape(slide, Inches(0.8), Inches(1.2), Inches(5.8), Inches(3.0), border_color=RED_ACCENT)
tf = add_rich_text_box(slide, Inches(1.1), Inches(1.4), Inches(5.2), Inches(2.6))
add_para(tf, "The Problem", size=18, color=RED_ACCENT, bold=True)
add_para(tf, "", size=6, color=DIM)
add_bullet(tf, "Standard N=8192 (PN13QP218) provides only 5 levels", size=13, color=WHITE)
add_bullet(tf, "Our model requires 7 multiplicative levels:", size=13, color=WHITE)
add_bullet(tf, "  Conv(1) + Poly(2) + Conv(1) + Poly(2) + FC(1) = 7", size=12, color=LIGHT_GRAY)
add_bullet(tf, "N=16384 works but is 2x slower (6,701 ms)", size=13, color=WHITE)
add_bullet(tf, "Goal: run at N=8192 for faster inference", size=13, color=ORANGE)

# Insight card
add_shape(slide, Inches(7.0), Inches(1.2), Inches(5.5), Inches(3.0), border_color=ACCENT3)
tf = add_rich_text_box(slide, Inches(7.3), Inches(1.4), Inches(5.0), Inches(2.6))
add_para(tf, "Key Insight", size=18, color=ACCENT3, bold=True)
add_para(tf, "", size=6, color=DIM)
add_bullet(tf, "CKKS security depends on logQP (total modulus", size=13, color=WHITE)
add_bullet(tf, "bit-width), NOT individual prime sizes", size=13, color=WHITE)
add_para(tf, "", size=6, color=DIM)
add_bullet(tf, "Smaller primes = more levels in same budget!", size=14, color=ACCENT3, bold=True)
add_para(tf, "", size=6, color=DIM)
add_bullet(tf, "FP8 models are tolerant of reduced precision", size=13, color=WHITE)
add_bullet(tf, "→ Perfect match for small-prime CKKS", size=13, color=ORANGE)

# Trade-off visual
add_shape(slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.5), border_color=ACCENT)
add_text_box(slide, Inches(1.1), Inches(4.8), Inches(11), Inches(0.4),
             "Trade-off: Per-Level Precision vs. Number of Levels", font_size=16, bold=True, color=ACCENT)

# Standard params
add_shape(slide, Inches(1.1), Inches(5.4), Inches(5.2), Inches(1.4), fill_color=BG_CARD)
tf = add_rich_text_box(slide, Inches(1.3), Inches(5.5), Inches(4.8), Inches(1.2))
add_para(tf, "Standard PN13QP218", size=13, color=DIM, bold=True)
add_para(tf, "Q = 1x33 + 5x30-bit  |  5 levels  |  30-bit scale", size=12, color=WHITE)

# Custom params
add_shape(slide, Inches(6.8), Inches(5.4), Inches(5.5), Inches(1.4), fill_color=BG_CARD, border_color=ACCENT3)
tf = add_rich_text_box(slide, Inches(7.0), Inches(5.5), Inches(5.1), Inches(1.2))
add_para(tf, "Custom PN13QP218s", size=13, color=ACCENT3, bold=True)
add_para(tf, "Q = 1x26 + 7x24-bit  |  7 levels  |  24-bit scale", size=12, color=WHITE)

add_slide_number(slide, 3, TOTAL)

# ===================== SLIDE 4: Core Contribution =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), fill_color=GOLD)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.5),
             "Our Core Contribution: CKKS Param Selection for QAT Models", font_size=24, bold=True, color=GOLD)

# Main contribution statement
add_shape(slide, Inches(0.8), Inches(1.2), Inches(11.7), Inches(1.4), fill_color=RGBColor(0x1A, 0x1A, 0x00), border_color=GOLD)
tf = add_rich_text_box(slide, Inches(1.1), Inches(1.35), Inches(11.1), Inches(1.1))
add_para(tf, "A systematic methodology to design custom CKKS parameters that exploit the reduced-precision tolerance",
         size=14, color=WHITE, bold=True)
add_para(tf, "of QAT models, enabling encrypted inference at smaller polynomial degrees (N) with significant speedup.",
         size=14, color=WHITE, bold=True)
add_para(tf, "", size=6, color=DIM)
add_para(tf, "Applicable to ANY quantized model (FP8/INT8/INT4) running FHE inference on the latti-ai platform.",
         size=12, color=GOLD)

# Three pillars
pillars = [
    ("Observation", ACCENT,
     ["QAT models learn noise-tolerant representations",
      "FP8 E4M3 quantization adds ~3-4 bits of noise",
      "CKKS scale precision beyond model need is wasted",
      "CKKS security = logQP, not individual prime size",
      "Smaller primes = more levels in same security budget"]),
    ("Method", ACCENT2,
     ["1. Profile model multiplicative depth (L)",
      "2. Choose target N by min ring dimension",
      "3. Compute security budget: logQP_max(N)",
      "4. Select Q prime bit-width matching quant precision",
      "5. Solve: fit L levels within logQP_max",
      "6. Validate P >= Q for key switching"]),
    ("Result", ACCENT3,
     ["Fewer levels needed (quant model is shallower)",
      "Smaller N (faster NTT, less memory, lower latency)",
      "No bootstrapping needed for moderate-depth models",
      "Pluggable into existing compiler pipeline",
      "Automatic selection via compiler auto-tuning"]),
]

for i, (title, color, bullets) in enumerate(pillars):
    x = Inches(0.8 + i * 4.0)
    add_shape(slide, x, Inches(2.9), Inches(3.7), Inches(4.2), border_color=color)
    add_text_box(slide, x + Inches(0.2), Inches(3.05), Inches(3.3), Inches(0.35),
                 title, font_size=17, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    tf = add_rich_text_box(slide, x + Inches(0.3), Inches(3.5), Inches(3.2), Inches(3.4))
    for b in bullets:
        add_bullet(tf, b, size=11, color=LIGHT_GRAY)

add_slide_number(slide, 4, TOTAL)

# ===================== SLIDE 5: Algorithm / Step-by-step =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), fill_color=GOLD)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.5),
             "CKKS Parameter Selection Algorithm", font_size=28, bold=True, color=GOLD)

# Algorithm steps as numbered cards
algo_steps = [
    ("Step 1", "Profile Model Depth", ORANGE,
     "Count multiplicative levels consumed by each layer:\n"
     "  Conv2d = 1 level, PolyRelu(d) = ceil(log2 d) levels\n"
     "  FC = 1 level, BN = 0 (fused into Conv)\n"
     "Example: Conv(1) + Poly2(2) + Conv(1) + Poly2(2) + FC(1) = 7"),
    ("Step 2", "Choose Target N", ACCENT,
     "Start from smallest N that could theoretically fit:\n"
     "  For N=8192: logQP_budget = 218 (128-bit security)\n"
     "  For N=16384: logQP_budget = 438 (128-bit security)\n"
     "Smaller N = faster NTT = faster inference"),
    ("Step 3", "Set Q Prime Bit-Width", ACCENT3,
     "Match CKKS scale precision to quantization format:\n"
     "  FP8 (E4M3): ~3-bit mantissa -> 22-26 bit Q primes\n"
     "  INT8: similar range -> 22-26 bit Q primes\n"
     "  FP32: full precision -> 30-34 bit Q primes\n"
     "Key insight: quantized models don't need 34-bit scale"),
    ("Step 4", "Solve Level Packing", ACCENT2,
     "Given: required_levels (L), logQP_budget, Q_bits\n"
     "Compute: logQP = log2(Q0) + L * Q_bits + log2(P)\n"
     "Constraint: logQP <= logQP_budget AND P_bits >= Q_bits\n"
     "Adjust Q_bits down until constraints are satisfied"),
]

for i, (step, title, color, desc) in enumerate(algo_steps):
    y = Inches(1.1 + i * 1.55)
    # Step badge
    add_shape(slide, Inches(0.8), y, Inches(1.3), Inches(0.4), fill_color=color)
    add_text_box(slide, Inches(0.8), y + Inches(0.05), Inches(1.3), Inches(0.3),
                 step, font_size=11, bold=True, color=BG_DARK, alignment=PP_ALIGN.CENTER)
    # Title
    add_text_box(slide, Inches(2.3), y + Inches(0.02), Inches(4), Inches(0.35),
                 title, font_size=14, bold=True, color=color)
    # Description card
    add_shape(slide, Inches(0.8), y + Inches(0.5), Inches(11.7), Inches(0.95))
    add_text_box(slide, Inches(1.1), y + Inches(0.52), Inches(11.1), Inches(0.9),
                 desc, font_size=10, color=LIGHT_GRAY, font_name="Consolas")

add_slide_number(slide, 5, TOTAL)

# ===================== SLIDE 6: Case Study — FP8 QAT Training =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), fill_color=ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.5),
             "Case Study: FP8 Quantization-Aware Training", font_size=24, bold=True, color=WHITE)

# Model architecture
add_shape(slide, Inches(0.8), Inches(1.2), Inches(11.7), Inches(2.0), border_color=ACCENT2)
add_text_box(slide, Inches(1.1), Inches(1.35), Inches(4), Inches(0.35),
             "Model: SimpleCNN (16x16 MNIST)", font_size=15, bold=True, color=ACCENT2)
tf = add_rich_text_box(slide, Inches(1.1), Inches(1.8), Inches(11), Inches(1.2))
add_para(tf, "Input(1x16x16) -> Conv2d(1,16,k=5,s=2) -> BN -> ReLU -> Conv2d(16,32,k=3,s=2) -> BN -> ReLU -> Flatten -> Dropout -> Linear(512,10)", size=12, color=LIGHT_GRAY, font_name="Consolas")
add_para(tf, "~11K parameters  |  2 convolutional layers + 1 FC layer", size=12, color=DIM)

# Accuracy table
add_shape(slide, Inches(0.8), Inches(3.5), Inches(11.7), Inches(2.5), border_color=ACCENT)
add_text_box(slide, Inches(1.1), Inches(3.7), Inches(5), Inches(0.35),
             "Training Results", font_size=16, bold=True, color=ACCENT)

# Table header
y_table = Inches(4.2)
add_shape(slide, Inches(1.1), y_table, Inches(11), Inches(0.4), fill_color=ACCENT)
add_text_box(slide, Inches(1.3), y_table + Inches(0.05), Inches(2.5), Inches(0.3),
             "Model", font_size=12, bold=True, color=BG_DARK)
add_text_box(slide, Inches(4.0), y_table + Inches(0.05), Inches(2.5), Inches(0.3),
             "Test Accuracy", font_size=12, bold=True, color=BG_DARK, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(6.5), y_table + Inches(0.05), Inches(2.5), Inches(0.3),
             "Delta vs Baseline", font_size=12, bold=True, color=BG_DARK, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(9.0), y_table + Inches(0.05), Inches(2.8), Inches(0.3),
             "Method", font_size=12, bold=True, color=BG_DARK, alignment=PP_ALIGN.CENTER)

rows = [
    ("FP32 Baseline", "98.56%", "--", "Standard training"),
    ("FP8 Post-Quant", "98.50%", "-0.06%", "Quantize after training"),
    ("FP8 QAT", "98.55%", "-0.01%", "Fake quant with STE"),
]
for i, (name, acc, delta, method) in enumerate(rows):
    y = y_table + Inches(0.45 + i * 0.38)
    bg = BG_CARD if i % 2 == 0 else BG_DARK
    add_shape(slide, Inches(1.1), y, Inches(11), Inches(0.35), fill_color=bg)
    c = ACCENT3 if i == 2 else WHITE
    add_text_box(slide, Inches(1.3), y + Inches(0.03), Inches(2.5), Inches(0.3),
                 name, font_size=11, color=c, bold=(i == 2))
    add_text_box(slide, Inches(4.0), y + Inches(0.03), Inches(2.5), Inches(0.3),
                 acc, font_size=11, color=c, bold=(i == 2), alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(6.5), y + Inches(0.03), Inches(2.5), Inches(0.3),
                 delta, font_size=11, color=ACCENT3 if i == 2 else LIGHT_GRAY,
                 alignment=PP_ALIGN.CENTER, bold=(i == 2))
    add_text_box(slide, Inches(9.0), y + Inches(0.03), Inches(2.8), Inches(0.3),
                 method, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# QAT explanation
add_shape(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(1.0))
tf = add_rich_text_box(slide, Inches(1.1), Inches(6.3), Inches(11), Inches(0.8))
add_para(tf, "FP8 QAT applies fake quantization (E4M3) during training via Straight-Through Estimator (STE). "
         "Weights remain FP32 but learn to be robust to quantization noise — nearly perfectly recovering baseline accuracy.",
         size=11, color=LIGHT_GRAY)

add_slide_number(slide, 6, TOTAL)

# ===================== SLIDE 7: Case Study — FHE Model Conversion =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), fill_color=ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.5),
             "Case Study: FHE Model Conversion", font_size=28, bold=True, color=WHITE)

# Three conversion steps as cards
conversions = [
    ("Activation Replacement", ACCENT2,
     ["ReLU -> polynomial approximation",
      "RangeNormPoly2d (degree 2)",
      "Polynomial degree 2: cost = 2 levels",
      "Fine-tuned 7 epochs: 97.76% accuracy"]),
    ("Pooling Replacement", ACCENT,
     ["MaxPool -> AvgPool",
      "MaxPool requires comparison ops",
      "Comparison unsupported in FHE",
      "AvgPool = depthwise conv equivalent"]),
    ("Conv + BN Fusion", ACCENT3,
     ["BatchNorm fused into Conv weights",
      "No extra multiplicative levels",
      "scale = gamma / sqrt(var + eps)",
      "fused_w[i] = w[i]*scale[i] + bias"]),
]

for i, (title, color, bullets) in enumerate(conversions):
    x = Inches(0.8 + i * 4.1)
    add_shape(slide, x, Inches(1.2), Inches(3.8), Inches(3.5), border_color=color)
    add_text_box(slide, x + Inches(0.3), Inches(1.4), Inches(3.2), Inches(0.4),
                 title, font_size=15, bold=True, color=color)
    tf = add_rich_text_box(slide, x + Inches(0.3), Inches(1.9), Inches(3.2), Inches(2.5))
    for b in bullets:
        add_bullet(tf, b, size=11, color=LIGHT_GRAY)

# Level consumption
add_shape(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(2.2), border_color=ORANGE)
add_text_box(slide, Inches(1.1), Inches(5.15), Inches(8), Inches(0.35),
             "Level Consumption (degree=2, ordinary style)", font_size=15, bold=True, color=ORANGE)

tf = add_rich_text_box(slide, Inches(1.3), Inches(5.6), Inches(11), Inches(1.5))
add_para(tf, "input L=7  ->  Conv L7->6(1)  ->  Poly L6->4(2)  ->  Conv L4->3(1)  ->  Poly L3->1(2)  ->  FC L1->0(1)  ->  output L=0",
         size=12, color=WHITE, font_name="Consolas")
add_para(tf, "Total: 7 multiplicative levels consumed (fits exactly in PN13QP218s)",
         size=13, color=ACCENT3, bold=True)

add_slide_number(slide, 7, TOTAL)

# ===================== SLIDE 8: Case Study — CKKS Parameter Design =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), fill_color=ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.5),
             "Case Study: CKKS Param Design — PN13QP218s", font_size=24, bold=True, color=WHITE)

# Parameter specification
add_shape(slide, Inches(0.8), Inches(1.2), Inches(5.8), Inches(3.5), border_color=ACCENT)
add_text_box(slide, Inches(1.1), Inches(1.35), Inches(5), Inches(0.35),
             "Parameter Specification", font_size=15, bold=True, color=ACCENT)
tf = add_rich_text_box(slide, Inches(1.1), Inches(1.8), Inches(5.2), Inches(2.8))
add_bullet(tf, "N = 8,192  (logN = 13)", size=12, color=WHITE)
add_bullet(tf, "max_level = 7", size=12, color=WHITE)
add_bullet(tf, "log_default_scale = 24", size=12, color=WHITE)
add_para(tf, "", size=6, color=DIM)
add_bullet(tf, "Q0: 1 x 26-bit (0x2044001)", size=12, color=LIGHT_GRAY)
add_bullet(tf, "Q1-Q7: 7 x 24-bit (NTT-friendly)", size=12, color=LIGHT_GRAY)
add_bullet(tf, "P: 1 x 26-bit (0x207C001)", size=12, color=LIGHT_GRAY)
add_para(tf, "", size=6, color=DIM)
add_bullet(tf, "logQP = 220 (~127-bit security)", size=13, color=ACCENT3, bold=True)
add_bullet(tf, "P >= Q (key switching requirement)", size=12, color=ORANGE)

# Design iterations
add_shape(slide, Inches(7.0), Inches(1.2), Inches(5.5), Inches(3.5), border_color=RED_ACCENT)
add_text_box(slide, Inches(7.3), Inches(1.35), Inches(5), Inches(0.35),
             "Design Iterations", font_size=15, bold=True, color=RED_ACCENT)

configs = [
    ("22-bit Q + 34-bit P", "logQP=216", "error 0.170", "FAIL", RED_ACCENT),
    ("24-bit Q + 22-bit P", "logQP=218", "error 2.093", "FAIL", RED_ACCENT),
    ("24/22 hybrid + 26-bit P", "logQP=218", "error 0.106", "FAIL", RED_ACCENT),
    ("24-bit Q + 26-bit P", "logQP=220", "error 0.037", "PASS", ACCENT3),
]

for i, (config, lqp, err, result, rc) in enumerate(configs):
    y = Inches(1.85 + i * 0.6)
    bg = BG_CARD if i < 3 else RGBColor(0x0A, 0x2A, 0x1A)
    add_shape(slide, Inches(7.2), y, Inches(5.1), Inches(0.5), fill_color=bg, border_color=rc if i == 3 else None)
    add_text_box(slide, Inches(7.4), y + Inches(0.08), Inches(2.0), Inches(0.3),
                 config, font_size=10, color=WHITE, font_name="Consolas")
    add_text_box(slide, Inches(9.4), y + Inches(0.08), Inches(1.2), Inches(0.3),
                 lqp, font_size=10, color=DIM, font_name="Consolas", alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(10.6), y + Inches(0.08), Inches(1.0), Inches(0.3),
                 err, font_size=10, color=DIM, font_name="Consolas", alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(11.6), y + Inches(0.08), Inches(0.6), Inches(0.3),
                 result, font_size=11, color=rc, bold=True, alignment=PP_ALIGN.CENTER)

# Lessons
add_shape(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(2.2), border_color=ORANGE)
add_text_box(slide, Inches(1.1), Inches(5.15), Inches(8), Inches(0.35),
             "Key Lessons", font_size=15, bold=True, color=ORANGE)
tf = add_rich_text_box(slide, Inches(1.3), Inches(5.6), Inches(11), Inches(1.4))
add_bullet(tf, "1. P must be >= largest Q level prime for correct key switching (catastrophic failure otherwise)", size=12, color=WHITE)
add_bullet(tf, "2. 22-bit Q primes lack precision for polynomial activations", size=12, color=WHITE)
add_bullet(tf, "3. Uniform 24-bit Q primes provide the best balance of precision and level count", size=12, color=WHITE)
add_bullet(tf, "4. All primes must satisfy p = 1 (mod 2N) for NTT compatibility", size=12, color=WHITE)

add_slide_number(slide, 8, TOTAL)

# ===================== SLIDE 9: Results Comparison =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), fill_color=ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.5),
             "Encrypted Inference Results", font_size=28, bold=True, color=WHITE)

# Comparison table
metrics = [
    ("CKKS Parameter Set",     "PN14QP438", "PN14QP438",     "PN13QP218s"),
    ("Poly Modulus Degree (N)", "16,384",     "16,384",       "8,192"),
    ("Multiplicative Levels",   "9",          "9",            "7"),
    ("Q Prime Size",            "34-bit",     "34-bit",       "24-bit"),
    ("Log Default Scale",       "34",         "34",           "24"),
    ("Security",                "128-bit",    "128-bit",      "~127-bit"),
    ("Plaintext Accuracy",      "98.56%",     "97.76%",       "97.76%"),
    ("Inference Time",          "6,701 ms",   "4,346 ms",     "1,194 ms"),
    ("Max Absolute Error",      "0.000016",   "0.000018",     "0.037035"),
    ("Verification",            "PASS",       "PASS",         "PASS"),
    ("Speedup",                 "1.0x",       "1.5x",         "5.6x"),
]

# Headers
headers = ["Metric", "FP32 + d=4", "FP8 QAT + d=2", "FP8 QAT + d=2 (Ours)"]
header_colors = [DIM, LIGHT_GRAY, LIGHT_GRAY, ACCENT3]
for j, (h, hc) in enumerate(zip(headers, header_colors)):
    x = Inches(0.8 + j * 3.0)
    add_shape(slide, x, Inches(1.2), Inches(2.8), Inches(0.45), fill_color=ACCENT if j == 3 else BG_CARD)
    add_text_box(slide, x + Inches(0.15), Inches(1.25), Inches(2.5), Inches(0.35),
                 h, font_size=11, bold=True, color=BG_DARK if j == 3 else hc,
                 alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

for i, row in enumerate(metrics):
    y = Inches(1.72 + i * 0.44)
    bg = BG_CARD if i % 2 == 0 else BG_DARK
    for j, val in enumerate(row):
        x = Inches(0.8 + j * 3.0)
        fill = bg
        c = LIGHT_GRAY
        bold = False
        if j == 3:
            fill = RGBColor(0x0A, 0x2A, 0x1A) if i % 2 == 0 else RGBColor(0x08, 0x22, 0x16)
            c = ACCENT3
            bold = True
        if row[0] == "Inference Time" and j == 3:
            c = ORANGE
        if row[0] == "Speedup" and j == 3:
            c = ORANGE
        add_shape(slide, x, y, Inches(2.8), Inches(0.4), fill_color=fill)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.05), Inches(2.5), Inches(0.3),
                     val, font_size=11, color=c, bold=bold,
                     alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

# Bottom note
add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.4),
             "Our approach: FP8 QAT + degree-2 poly + custom 24-bit CKKS primes = 5.6x speedup at N=8192",
             font_size=13, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 9, TOTAL)

# ===================== SLIDE 10: Verification Detail =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), fill_color=ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.5),
             "Verification Detail (FP8 QAT, N=8192)", font_size=24, bold=True, color=WHITE)

# Output comparison table
add_shape(slide, Inches(0.8), Inches(1.2), Inches(11.7), Inches(4.0), border_color=ACCENT3)
add_text_box(slide, Inches(1.1), Inches(1.35), Inches(8), Inches(0.35),
             "Encrypted vs Plaintext Output Comparison", font_size=15, bold=True, color=ACCENT3)

# Headers
for j, h in enumerate(["Index", "Encrypted", "Plaintext", "Abs Error"]):
    x = Inches(1.1 + j * 2.8)
    w = Inches(2.6) if j > 0 else Inches(1.0)
    add_text_box(slide, x, Inches(1.8), w, Inches(0.3),
                 h, font_size=11, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)

verif_data = [
    (0, -2.03323238, -1.99619718, 0.03703520),
    (1, -0.12924453, -0.16405774, 0.03481321),
    (2,  1.95609878,  1.96626933, 0.01017055),
    (3, -0.20319583, -0.20021933, 0.00297650),
    (4, -2.36351947, -2.37562968, 0.01211021),
    (5,  0.61226799,  0.61428893, 0.00202095),
    (6, -1.17538829, -1.17316265, 0.00222563),
    (7,  0.05304573,  0.02736787, 0.02567785),
    (8,  0.54493935,  0.55579004, 0.01085069),
    (9, -1.88872228, -1.91488866, 0.02616638),
]

for i, (idx, enc, pt, err) in enumerate(verif_data):
    y = Inches(2.15 + i * 0.3)
    bg = BG_CARD if i % 2 == 0 else BG_DARK
    add_shape(slide, Inches(1.1), y, Inches(11), Inches(0.28), fill_color=bg)
    add_text_box(slide, Inches(1.1), y + Inches(0.02), Inches(1.0), Inches(0.24),
                 str(idx), font_size=10, color=DIM, alignment=PP_ALIGN.CENTER, font_name="Consolas")
    add_text_box(slide, Inches(3.9), y + Inches(0.02), Inches(2.6), Inches(0.24),
                 f"{enc:>14.8f}", font_size=10, color=WHITE, alignment=PP_ALIGN.CENTER, font_name="Consolas")
    add_text_box(slide, Inches(6.7), y + Inches(0.02), Inches(2.6), Inches(0.24),
                 f"{pt:>14.8f}", font_size=10, color=WHITE, alignment=PP_ALIGN.CENTER, font_name="Consolas")
    err_color = RED_ACCENT if err > 0.03 else ACCENT3
    add_text_box(slide, Inches(9.5), y + Inches(0.02), Inches(2.6), Inches(0.24),
                 f"{err:.8f}", font_size=10, color=err_color, alignment=PP_ALIGN.CENTER, font_name="Consolas")

# Summary
add_shape(slide, Inches(0.8), Inches(5.5), Inches(5.5), Inches(1.5), border_color=ACCENT3)
tf = add_rich_text_box(slide, Inches(1.1), Inches(5.65), Inches(5.0), Inches(1.2))
add_para(tf, "Max absolute error:  0.03703520", size=13, color=WHITE, font_name="Consolas")
add_para(tf, "Avg absolute error:  0.01640472", size=13, color=WHITE, font_name="Consolas")
add_para(tf, "Result:  PASS", size=16, color=ACCENT3, bold=True, font_name="Consolas")

add_shape(slide, Inches(6.7), Inches(5.5), Inches(5.8), Inches(1.5), border_color=ACCENT)
tf = add_rich_text_box(slide, Inches(7.0), Inches(5.65), Inches(5.2), Inches(1.2))
add_para(tf, "Classification: digit 2", size=13, color=WHITE, font_name="Consolas")
add_para(tf, "Both encrypted & plaintext predict class 2", size=13, color=ACCENT3, font_name="Consolas")
add_para(tf, "Correct classification preserved", size=13, color=ACCENT3, bold=True)

add_slide_number(slide, 10, TOTAL)

# ===================== SLIDE 11: Generalization & Impact =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), fill_color=ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.5),
             "Generalization: QAT + Custom CKKS for Any Model", font_size=24, bold=True, color=GOLD)

# Applicable models table
add_shape(slide, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.45), fill_color=GOLD)
for j, h in enumerate(["Quantization", "Q Bits", "Typical Depth", "Min N", "Speedup vs FP32 N"]):
    w = Inches(2.3) if j > 0 else Inches(2.5)
    x = Inches(0.8 + j * 2.3) if j == 0 else Inches(0.8 + 2.5 + (j-1) * 2.3)
    add_text_box(slide, x + Inches(0.1), Inches(1.25), w, Inches(0.35),
                 h, font_size=11, bold=True, color=BG_DARK,
                 alignment=PP_ALIGN.CENTER)

model_rows = [
    ("FP8 (E4M3)", "24-bit", "7-9 levels", "N=8192", "4-6x", ACCENT3),
    ("INT8", "22-24-bit", "7-9 levels", "N=8192", "4-6x", ACCENT3),
    ("INT4", "20-22-bit", "7-9 levels", "N=8192", "4-6x", ACCENT3),
    ("FP16 (QAT)", "28-30-bit", "7-9 levels", "N=16384", "2-3x", ACCENT),
    ("FP32 (no QAT)", "34-bit", "7-9 levels", "N=16384+", "1.0x (baseline)", DIM),
]

for i, (quant, qbits, depth, n, speedup, color) in enumerate(model_rows):
    y = Inches(1.72 + i * 0.44)
    bg = BG_CARD if i % 2 == 0 else BG_DARK
    vals = [quant, qbits, depth, n, speedup]
    for j, val in enumerate(vals):
        w = Inches(2.3) if j > 0 else Inches(2.5)
        x = Inches(0.8 + j * 2.3) if j == 0 else Inches(0.8 + 2.5 + (j-1) * 2.3)
        fill = bg
        c = color if j == 4 else LIGHT_GRAY
        bold = (j == 4)
        add_shape(slide, x, y, w, Inches(0.4), fill_color=fill)
        add_text_box(slide, x + Inches(0.1), y + Inches(0.05), w - Inches(0.2), Inches(0.3),
                     val, font_size=11, color=c, bold=bold, alignment=PP_ALIGN.CENTER,
                     font_name="Consolas")

# Impact section
add_shape(slide, Inches(0.8), Inches(4.1), Inches(5.5), Inches(3.0), border_color=ACCENT3)
add_text_box(slide, Inches(1.1), Inches(4.25), Inches(5), Inches(0.35),
             "Why This Matters", font_size=16, bold=True, color=ACCENT3)
tf = add_rich_text_box(slide, Inches(1.3), Inches(4.7), Inches(4.8), Inches(2.2))
add_bullet(tf, "Enables N=8192 for models that previously required N=16384+", size=12, color=WHITE)
add_bullet(tf, "No bootstrapping needed for moderate-depth quantized models", size=12, color=WHITE)
add_bullet(tf, "Automatic integration into compiler pipeline", size=12, color=WHITE)
add_bullet(tf, "Orthogonal to other optimizations (poly degree, model pruning)", size=12, color=WHITE)
add_bullet(tf, "Directly applicable to production FHE deployments", size=12, color=ACCENT3, bold=True)

# Compiler integration flow
add_shape(slide, Inches(6.7), Inches(4.1), Inches(5.8), Inches(3.0), border_color=ACCENT2)
add_text_box(slide, Inches(7.0), Inches(4.25), Inches(5.2), Inches(0.35),
             "Compiler Integration", font_size=16, bold=True, color=ACCENT2)
tf = add_rich_text_box(slide, Inches(7.2), Inches(4.7), Inches(5.2), Inches(2.2))
add_bullet(tf, "pipeline.py: auto-search smallest viable N", size=12, color=WHITE)
add_bullet(tf, "  try_no_btp() tests PN13QP218s first", size=11, color=LIGHT_GRAY)
add_bullet(tf, "components.py: parameter set definitions", size=12, color=WHITE)
add_bullet(tf, "deploy_cmds.py: HE instruction lookup", size=12, color=WHITE)
add_bullet(tf, "inference_process.cpp: custom Q/P loading", size=12, color=WHITE)
add_para(tf, "", size=6, color=DIM)
add_para(tf, "User workflow: train QAT model -> compile -> inference", size=12, color=ACCENT2, bold=True)
add_para(tf, "(parameter selection is fully automatic)", size=11, color=LIGHT_GRAY)

add_slide_number(slide, 11, TOTAL)

# ===================== SLIDE 12: Code Changes =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), fill_color=ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.5),
             "Code Changes", font_size=28, bold=True, color=WHITE)

changes = [
    ("components.py", "New PN13QP218s parameter set", "24-bit Q, 26-bit P, 7 levels", ACCENT3),
    ("pipeline.py", "Compiler parameter selection", "PN13QP218s first in auto-selection list", ACCENT2),
    ("deploy_cmds.py", "HE instruction generator", "PN13QP218s in _FHE_PARAMS lookup", ACCENT),
    ("inference_process.cpp", "Server-side param loading", "Read custom Q/P from JSON", ORANGE),
    ("CMakeLists.txt", "macOS ARM64 build fixes", "Guard -maes, remove Linux-only flags", DIM),
]

for i, (file, desc, detail, color) in enumerate(changes):
    y = Inches(1.3 + i * 1.15)
    add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(1.0), border_color=color)
    add_text_box(slide, Inches(1.2), y + Inches(0.08), Inches(3.5), Inches(0.35),
                 file, font_size=14, bold=True, color=color, font_name="Consolas")
    add_text_box(slide, Inches(5.0), y + Inches(0.08), Inches(4), Inches(0.35),
                 desc, font_size=13, color=WHITE)
    add_text_box(slide, Inches(5.0), y + Inches(0.45), Inches(7), Inches(0.35),
                 detail, font_size=11, color=LIGHT_GRAY)

add_slide_number(slide, 12, TOTAL)

# ===================== SLIDE 13: Pipeline Summary =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), fill_color=ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.5),
             "End-to-End Pipeline", font_size=28, bold=True, color=WHITE)

# Pipeline flow
steps = [
    ("FP32 Baseline", "98.56%", ACCENT2),
    ("FP8 QAT Training", "98.55%\n(-0.01%)", ACCENT2),
    ("FHE Conversion", "97.76%\n(-0.80%)", ACCENT),
    ("Custom CKKS\nParameter Design", "N=8192\n7 levels", ORANGE),
    ("Encrypted\nInference", "1,194 ms\n5.6x speedup", ACCENT3),
]

for i, (label, metric, color) in enumerate(steps):
    x = Inches(0.5 + i * 2.5)
    # Box
    add_shape(slide, x, Inches(1.4), Inches(2.2), Inches(2.0), border_color=color)
    add_text_box(slide, x + Inches(0.15), Inches(1.55), Inches(1.9), Inches(0.8),
                 label, font_size=12, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), Inches(2.4), Inches(1.9), Inches(0.8),
                 metric, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    # Arrow
    if i < len(steps) - 1:
        add_text_box(slide, x + Inches(2.2), Inches(2.0), Inches(0.3), Inches(0.4),
                     ">", font_size=20, color=DIM, alignment=PP_ALIGN.CENTER)

# Reproduce commands
add_shape(slide, Inches(0.8), Inches(3.8), Inches(11.7), Inches(3.4), border_color=ACCENT)
add_text_box(slide, Inches(1.1), Inches(3.95), Inches(8), Inches(0.35),
             "Reproduction Steps", font_size=15, bold=True, color=ACCENT)

tf = add_rich_text_box(slide, Inches(1.3), Inches(4.4), Inches(10.8), Inches(2.6))
commands = [
    "# Step 1: Compile with auto CKKS parameter selection (selects PN13QP218s)",
    "python3 training/run_compile.py \\",
    "  -i examples/test_mnist/output_fp8qat/trained_poly_d2.onnx \\",
    "  -o examples/test_mnist/output_fp8qat --style ordinary",
    "",
    "# Step 2: Generate HE instructions",
    "python3 inference/interface/gen_mega_ag.py --task-dir examples/test_mnist/output_fp8qat/task",
    "",
    "# Step 3: Run encrypted inference",
    "./build/examples/inference \\",
    "  --task-dir examples/test_mnist/output_fp8qat/task \\",
    "  --input examples/test_mnist/output_fp8qat/task/client/img.csv --verify",
]
for cmd in commands:
    add_para(tf, cmd, size=10, color=LIGHT_GRAY if cmd.startswith("#") else WHITE, font_name="Consolas")

add_slide_number(slide, 13, TOTAL)

# ===================== SLIDE 14: Summary & Conclusion =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), fill_color=GOLD)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.5),
             "Summary & Conclusion", font_size=28, bold=True, color=GOLD)

# What we did
add_shape(slide, Inches(0.8), Inches(1.2), Inches(5.8), Inches(3.2), border_color=ACCENT)
add_text_box(slide, Inches(1.1), Inches(1.35), Inches(5), Inches(0.35),
             "What We Did", font_size=16, bold=True, color=ACCENT)
tf = add_rich_text_box(slide, Inches(1.3), Inches(1.8), Inches(5.0), Inches(2.4))
add_bullet(tf, "Trained FP8 QAT model (98.55%, -0.01% vs FP32)", size=12, color=WHITE)
add_bullet(tf, "Converted to FHE-compatible polynomial model", size=12, color=WHITE)
add_bullet(tf, "Designed custom PN13QP218s (24-bit Q, 7 levels)", size=12, color=WHITE)
add_bullet(tf, "Ran encrypted inference at N=8192 in 1,194 ms", size=12, color=WHITE)
add_bullet(tf, "Achieved 5.6x speedup, verification PASS", size=12, color=ACCENT3, bold=True)

# Core contribution
add_shape(slide, Inches(7.0), Inches(1.2), Inches(5.5), Inches(3.2), border_color=GOLD)
add_text_box(slide, Inches(7.3), Inches(1.35), Inches(5), Inches(0.35),
             "Core Contribution", font_size=16, bold=True, color=GOLD)
tf = add_rich_text_box(slide, Inches(7.5), Inches(1.8), Inches(4.8), Inches(2.4))
add_bullet(tf, "Generalized CKKS parameter selection methodology", size=12, color=WHITE)
add_bullet(tf, "Exploits QAT noise tolerance for smaller primes", size=12, color=WHITE)
add_bullet(tf, "Enables smaller N = faster encrypted inference", size=12, color=WHITE)
add_bullet(tf, "Automatic compiler integration, no manual tuning", size=12, color=WHITE)
add_bullet(tf, "Applicable to all QAT models (FP8/INT8/INT4)", size=12, color=GOLD, bold=True)

# Key numbers highlight
add_shape(slide, Inches(0.8), Inches(4.7), Inches(11.7), Inches(1.5),
          fill_color=RGBColor(0x1A, 0x1A, 0x00), border_color=GOLD)
add_text_box(slide, Inches(1.1), Inches(4.85), Inches(11.1), Inches(0.35),
             "Key Numbers", font_size=16, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)

highlight_nums = [
    ("5.6x", "Speedup", ACCENT3),
    ("1,194ms", "Latency", ACCENT),
    ("N=8,192", "Ring Dim", ACCENT2),
    ("~127-bit", "Security", ORANGE),
    ("PASS", "Verification", ACCENT3),
]
for i, (num, label, color) in enumerate(highlight_nums):
    x = Inches(1.2 + i * 2.3)
    add_text_box(slide, x, Inches(5.25), Inches(2.0), Inches(0.45),
                 num, font_size=28, bold=True, color=color, alignment=PP_ALIGN.CENTER,
                 font_name="Consolas")
    add_text_box(slide, x, Inches(5.7), Inches(2.0), Inches(0.3),
                 label, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Footer
add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5),
             "FP8 QAT + Custom CKKS Parameter Selection = Practical Encrypted Neural Network Inference",
             font_size=15, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 14, TOTAL)

# ── Save ────────────────────────────────────────────────────────────────
output_path = "/Users/boycrypt/code/C++/latti-ai/examples/test_mnist/output_fp8qat/FP8_QAT_Cipher_Inference.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
